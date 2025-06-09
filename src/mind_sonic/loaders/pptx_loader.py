#!/usr/bin/env python
"""
PowerPoint Loader for EmbedChain

This module provides a custom loader for PowerPoint (PPTX) files.
"""
from typing import Dict, Any, List
import os
from pptx import Presentation
from embedchain.loaders.base_loader import BaseLoader

# Constants for metadata
FILE_TYPE = "pptx"


class PowerPointLoader(BaseLoader):
    """Custom loader for PowerPoint (.pptx) files.
    
    This loader extracts text content from PowerPoint presentations and formats it
    in a way that's compatible with the embedchain BaseChunker. It processes each slide
    in the presentation, extracts text from all shapes, and combines them into a single
    document with slide numbers for context.
    
    The loader returns data in the specific format expected by BaseChunker:
    {
        "data": [{
            "content": "extracted text content",
            "meta_data": {
                "source": "file path",
                "file_type": "pptx",
                "slide_count": number_of_slides,
                "url": "file path"  # BaseChunker expects a url in metadata
            }
        }],
        "doc_id": "unique document identifier"
    }
    
    Usage:
        loader = PowerPointLoader()
        data = loader.load_data("path/to/presentation.pptx")
    """
    
    def __init__(self):
        """Initialize the PowerPoint loader."""
        super().__init__()

    def load_data(self, source: str, **kwargs) -> Dict[str, Any]:
        """Load data from a PowerPoint file.
        
        This method extracts text content from a PowerPoint presentation file.
        It processes each slide, extracts text from all shapes on the slide,
        and formats the output with slide numbers for better context preservation.
        
        The method returns data in the specific format required by the embedchain
        BaseChunker, which expects a dictionary with 'data' and 'doc_id' keys.
        
        Args:
            source: Path to the PowerPoint file
            **kwargs: Additional keyword arguments that might be passed by the caller
                      (required for compatibility with BaseChunker)
            
        Returns:
            Dictionary containing the extracted text and metadata in the format expected by BaseChunker:
            {
                "data": [{
                    "content": "extracted text",
                    "meta_data": {...}
                }],
                "doc_id": "unique identifier"
            }
            
        Raises:
            FileNotFoundError: If the specified file does not exist
        """
        if not os.path.exists(source):
            raise FileNotFoundError(f"File {source} not found")
            
        # Extract text from PowerPoint slides
        presentation = Presentation(source)
        all_text = []
        
        # Process each slide
        for i, slide in enumerate(presentation.slides):
            slide_text = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            # Add slide number and text to the result
            if slide_text:
                slide_content = f"Slide {i+1}:\n" + "\n".join(slide_text)
                all_text.append(slide_content)
        
        # Combine all text from slides
        full_text = "\n\n".join(all_text)
        
        # Generate a document ID based on the source file path
        import hashlib
        doc_id = hashlib.sha256(source.encode()).hexdigest()
        
        # Return the extracted text with metadata in the format expected by BaseChunker
        # Note: BaseChunker expects a dictionary with 'data' and 'doc_id' keys
        # The 'data' key should contain a list of dictionaries with 'content' and 'meta_data'
        # The 'meta_data' must include a 'url' field for BaseChunker to work properly
        return {
            "data": [{
                "content": full_text,
                "meta_data": {
                    "source": source,
                    "file_type": FILE_TYPE,
                    "slide_count": len(presentation.slides),
                    "url": source  # BaseChunker expects a url in metadata
                }
            }],
            "doc_id": doc_id
        }
